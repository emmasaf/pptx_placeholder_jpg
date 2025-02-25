import json

from pptx import Presentation
from pptx.util import Inches
from minio import Minio, S3Error
import os
import tempfile
import sys
import subprocess
from qrcode import QRCode
from pdf2image import convert_from_path
import io

if len(sys.argv) < 3:
    print("Usage: python3 generate_jpg.py <pptx_storage_key> <minio_data_json> ")
    sys.exit(1)

key = sys.argv[1]
minio_data_json = sys.argv[2]

minio_data = json.loads(minio_data_json)

if not isinstance(minio_data, dict):
    print("Error: 'minio_data_json' should be a dictionary.")
    sys.exit(1)

replacements_data = {
    "{{fullName}}": "Alice Smith",
    "{{status}}": "successfully completed",
    "{{date}}": "March 2025",
    "{{courseName}}": "Advanced Robotics",
    "{{trainerFullName}}": "Michael Johnson",
    "{{commercialDirector}}": "Sophia Reynolds",
    "{{uuid}}": "f47ac10b-58cc-4372-a567-0e02b2c3d479",
    "{{min_uuid}}": "944a7ca05e3244af618909",
    "{{qr_link}}": "https://example.com/qr-code"
}

client = Minio(
    minio_data['endpoint'],
    access_key=minio_data['access_key'],
    secret_key=minio_data['secret_key'],
    secure=False,
)

pdf_id = replacements_data["{{uuid}}"]
def convert_pdf_to_images(pdf_file_path, key_uuid):
    """Convert PDF to image and upload to MinIO."""
    try:
        if not os.path.exists(pdf_file_path):
            print(f"Error: PDF file {pdf_file_path} not found.")
            return

        images = convert_from_path(pdf_file_path)
        image_byte_stream = io.BytesIO()
        images[0].save(image_byte_stream, 'JPEG')
        image_byte_stream.seek(0)

        minio_key = f"certificates/jpg/test/{key_uuid}.jpg"
        client.put_object(
            "public", minio_key, image_byte_stream, len(image_byte_stream.getvalue())
        )
        print(f"Image uploaded to MinIO at {minio_key}")

    except S3Error as e:
        print(f"MinIO Error: {e}")
    except Exception as e:
        print(f"Error converting PDF to images: {e}")


def generate_pdf_and_upload(template_path, replacements):
    """Generate a PPTX, convert to PDF, and upload to MinIO."""
    prs = Presentation(template_path)

    # Generate QR code
    qr = QRCode(version=1, box_size=4, border=1)
    qr_data = replacements_data["{{qr_link}}"]
    qr.add_data(qr_data)
    qr.make(fit=True)

    qr_path = f"/tmp/{pdf_id}.png"
    qr.make_image(fill_color="black", back_color="white").save(qr_path)

    # Replace placeholders in PPTX
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == 'qr':
                left_in = shape.left / 914400
                top_in = shape.top / 914400
                width_in = shape.width / 914400

                slide.shapes._spTree.remove(shape._element)
                prs.slides[0].shapes.add_picture(qr_path, Inches(left_in), Inches(top_in), Inches(width_in),
                                                 Inches(width_in))

            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in paragraph.runs)
                    for placeholder, value in replacements.items():
                        if placeholder in full_text:
                            full_text = full_text.replace(placeholder, value)

                    if paragraph.runs:
                        paragraph.runs[0].text = full_text
                        for run in paragraph.runs[1:]:
                            run.text = ""

    # Save PPTX
    pptx_path = f"/tmp/{pdf_id}.pptx"
    prs.save(pptx_path)

    # Convert PPTX to PDF using LibreOffice (instead of unoconv)
    pdf_path = f"/tmp/{pdf_id}.pdf"
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", "/tmp"],
            check=True
        )
    except subprocess.CalledProcessError as e:
        print(f"Error converting PPTX to PDF: {e}")
        return

    # Ensure PDF file exists before proceeding
    if not os.path.exists(pdf_path):
        print(f"Error: Failed to generate PDF {pdf_path}")
        return

    destination_file = f"certificates/pdf/{pdf_id}.pdf"
    convert_pdf_to_images(pdf_path, pdf_id)

    os.remove(pptx_path)
    os.remove(pdf_path)
    os.remove(qr_path)

    print(f"Certificate successfully uploaded to bucket 'public'")
    print(f"Storage object key: {pdf_id}")

    return pdf_id


# Download PPTX from MinIO
bucket_name = "public"
object_name = key

tmp_file_path = os.path.join(tempfile.gettempdir(), os.path.basename(key))

if not os.path.exists(tmp_file_path):
    try:
        client.fget_object(bucket_name, object_name, tmp_file_path)
    except S3Error as e:
        print(f"Error downloading file: {e}")
        sys.exit(1)
else:
    print(f"'{tmp_file_path}' already exists.")

generate_pdf_and_upload(tmp_file_path, replacements_data)
os.remove(tmp_file_path)
sys.exit(0)
